from flask import Flask, request, jsonify, send_file
from werkzeug.utils import secure_filename
import os
from pptx import Presentation
import pdfkit
import tempfile
from flask_cors import CORS  # Added for CORS support

app = Flask(__name__)
CORS(app)  
app.config['UPLOAD_FOLDER'] = tempfile.gettempdir()
app.config['ALLOWED_EXTENSIONS'] = {'pptx'}
app.config['MAX_CONTENT_LENGTH'] = 10 * 1024 * 1024 

PDFKIT_CONFIG = pdfkit.configuration(wkhtmltopdf=r"C:\Program Files\wkhtmltopdf\bin\wkhtmltopdf.exe") # path to wkhtmltopdf executable

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    slides_text = []
    
    for i, slide in enumerate(prs.slides):
        slide_text = f"Slide {i+1}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        slides_text.append(slide_text)
    
    return slides_text

@app.route('/convert', methods=['POST'])
def convert_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    
    if not allowed_file(file.filename):
        return jsonify({'error': 'Invalid file type. Only PPTX files are allowed.'}), 400
    
    try:
        filename = secure_filename(file.filename)
        pptx_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(pptx_path)
        
        # Extract text from PPTX
        extracted_text = extract_text_from_pptx(pptx_path)
        
        # Create PDF path
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], filename + '.pdf')
        
        # Create HTML content
        html_content = "<html><body>"
        for i, slide_text in enumerate(extracted_text):
            html_content += f"<h2>Slide {i+1}</h2>"
            html_content += f"<pre>{slide_text}</pre><hr>"
        html_content += "</body></html>"
        
        # Convert HTML to PDF
        pdfkit.from_string(html_content, pdf_path, configuration=PDFKIT_CONFIG)
        
        return jsonify({
            'message': 'Conversion successful',
            'extracted_text': extracted_text,
            'pdf_url': f'/download/{filename}.pdf'
        })
        
    except Exception as e:
        return jsonify({'error': f'Conversion failed: {str(e)}'}), 500
    finally:
        if 'pptx_path' in locals() and os.path.exists(pptx_path):
            os.remove(pptx_path)

@app.route('/download/<filename>', methods=['GET'])
def download_file(filename):
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    if not os.path.exists(file_path):
        return jsonify({'error': 'File not found'}), 404
    
    try:
        return send_file(file_path, as_attachment=True)
    except Exception as e:
        return jsonify({'error': f'Download failed: {str(e)}'}), 500

if __name__ == '__main__':
    app.run(debug=True)